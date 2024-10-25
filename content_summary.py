from pptx import Presentation  # pptx 모듈에서 Presentation 클래스를 import

# PowerPoint 파일 열기
ppt_path = "C:/Users/Um/Desktop/Aichatbot/BCC.pptx"  # 첨부된 파일 경로로 수정하세요
presentation = Presentation(ppt_path)

# 슬라이드에서 텍스트 추출하기
content = ""
for slide in presentation.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            content += shape.text + "\n"

# 추출된 내용을 파일로 저장 (선택사항)
with open("content_summary.txt", "w", encoding="utf-8") as f:
    f.write(content)

print("PowerPoint content extracted successfully.")
